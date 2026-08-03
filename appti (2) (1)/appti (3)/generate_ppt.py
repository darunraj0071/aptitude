import os
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Image preparation: Create watermark from logo.png
logo_path = r"c:\Users\sanja\Downloads\appti (2) (1)\appti (3)\appti\images\logo.png"
faded_logo_path = "faded_logo.png"

try:
    logo = Image.open(logo_path).convert("RGBA")
    r, g, b, alpha = logo.split()
    # Apply extremely light opacity for the watermark (approx 8%)
    alpha_faded = ImageEnhance.Brightness(alpha).enhance(0.08)
    faded_logo = Image.merge("RGBA", (r, g, b, alpha_faded))
    faded_logo.save(faded_logo_path)
    print("Faded logo watermark created successfully.")
except Exception as e:
    print(f"Error creating faded logo: {e}")
    # Fallback to copy if PIL fails
    if os.path.exists(logo_path):
        import shutil
        shutil.copy(logo_path, faded_logo_path)

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Constants
BG_COLOR = RGBColor(15, 14, 23)        # #0F0E17 (Deep Dark Slate)
TEXT_WHITE = RGBColor(255, 255, 255)  # #FFFFFF
TEXT_SILVER = RGBColor(148, 161, 178) # #94A1B2 (Silver Gray)
ACCENT_PURPLE = RGBColor(127, 90, 240) # #7F5AF0 (Vibrant Purple)
ACCENT_GREEN = RGBColor(44, 182, 125)  # #2CB67D (Secondary Mint Green)
PANEL_BG = RGBColor(31, 30, 39)        # #1F1E27 (Darker Gray Panel)
BORDER_COLOR = RGBColor(60, 60, 80)

# Helper function to add slide title
def add_slide_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.93), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = TEXT_WHITE
    
    # Underline line
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.2), Inches(2.0), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ACCENT_PURPLE
    underline.line.fill.background()

# Helper function to apply background styling
def apply_slide_bg(slide, use_watermark=True):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if use_watermark and os.path.exists(faded_logo_path):
        # Center the watermark
        width_in = 5.0
        height_in = 5.0
        left_in = (13.333 - width_in) / 2
        top_in = (7.5 - height_in) / 2
        slide.shapes.add_picture(faded_logo_path, Inches(left_in), Inches(top_in), width=Inches(width_in), height=Inches(height_in))

# Helper to create bullet points standard slide
def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, use_watermark=True)
    add_slide_title(slide, title)
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.93), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(12)
        p.level = 0
        
        if isinstance(item, tuple):
            lead, body = item
            run = p.add_run()
            run.text = "• " + lead + ": "
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = ACCENT_PURPLE
            run.font.name = "Segoe UI"
            
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = TEXT_WHITE
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(16)
            run.font.color.rgb = TEXT_WHITE
            run.font.name = "Segoe UI"
    return slide

# Helper to create two column layout slides
def add_two_column_slide(prs, title, col1_title, col1_bullets, col2_title, col2_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, use_watermark=True)
    add_slide_title(slide, title)
    
    # Column 1
    tx1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.7), Inches(5.2))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = col1_title
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GREEN
    p1.font.name = "Segoe UI"
    p1.space_after = Pt(14)
    
    for item in col1_bullets:
        p = tf1.add_paragraph()
        p.space_after = Pt(10)
        if isinstance(item, tuple):
            lead, body = item
            run = p.add_run()
            run.text = "• " + lead + ": "
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = ACCENT_PURPLE
            run.font.name = "Segoe UI"
            
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(14)
            run2.font.color.rgb = TEXT_WHITE
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT_WHITE
            run.font.name = "Segoe UI"
            
    # Column 2
    tx2 = slide.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.7), Inches(5.2))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = col2_title
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.name = "Segoe UI"
    p2.space_after = Pt(14)
    
    for item in col2_bullets:
        p = tf2.add_paragraph()
        p.space_after = Pt(10)
        if isinstance(item, tuple):
            lead, body = item
            run = p.add_run()
            run.text = "• " + lead + ": "
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = ACCENT_PURPLE
            run.font.name = "Segoe UI"
            
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(14)
            run2.font.color.rgb = TEXT_WHITE
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT_WHITE
            run.font.name = "Segoe UI"
    return slide

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide1, use_watermark=False)
if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(5.41), Inches(0.8), width=Inches(2.5), height=Inches(2.5))

tx1 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.33), Inches(3.0))
tf1 = tx1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "VETRIPATH"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(8)

p2 = tf1.add_paragraph()
p2.text = "Complete Career Prep Portal & Android App"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = ACCENT_PURPLE
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(24)

p3 = tf1.add_paragraph()
p3.text = "An offline-first hybrid platform for placement, aptitude, and coding mastery"
p3.font.size = Pt(15)
p3.font.color.rgb = TEXT_SILVER
p3.font.name = "Segoe UI"
p3.alignment = PP_ALIGN.CENTER

# Slide 2: Abstract
abstract_bullets = [
    ("Concept & Purpose", "VetriPath is a self-contained, hybrid education system engineered to assist students preparing for campus placements, technical interviews, and exams."),
    ("Offline-First Engineering", "Solves connectivity barriers by packaging all application assets (HTML/CSS/JS) locally within an Android wrapper container, operating entirely offline without API latency."),
    ("Interactive Multi-Modal Training", "Hosts Quantitative Aptitude, Logical Reasoning, and Verbal Ability practice modules along with an integrated coding workstation."),
    ("Security Container", "Leverages Jetpack Compose and biometrics to provide private locking mechanism, preventing unauthorized access to user progress files."),
    ("Exam Safeguards", "Implements robust, client-side cheating protection protocols specifically tailored to replicate strict, closed-environment corporate assessments.")
]
add_bullet_slide(prs, "Project Abstract", abstract_bullets)

# Slide 3: Requirements
sw_bullets = [
    ("Operating System", "Android 8.0 (API 26) or higher for Mobile App; modern web browsers supporting HTML5/ES6 for Web Interface"),
    ("Core Languages", "Kotlin 1.9, Java (Android SDK), JavaScript (ES6+), HTML5, CSS3"),
    ("Frameworks/SDKs", "Jetpack Compose, Android Jetpack Biometrics SDK, WebView Engine"),
    ("Tools & Servers", "Android Studio Hedgehog, VS Code, Gradle Build Tool, local development python servers")
]
hw_bullets = [
    ("Development Unit", "Intel Core i5/AMD Ryzen 5 or higher, 8 GB RAM (16 GB recommended), 10 GB storage"),
    ("Testing Device", "Android Smartphone with Biometric Sensor, Touch screen interface (minimum 2 GB RAM)"),
    ("Candidate Hardware", "Core i3/equivalents, 4 GB RAM, minimum screen resolution 1366x768 pixels")
]
add_two_column_slide(prs, "Software & Hardware Requirements", "Software Environment Setup", sw_bullets, "Hardware Configuration Setup", hw_bullets)

# Slide 4: Existing System
existing_bullets = [
    ("Network Dependencies", "Most assessment prep platforms (e.g. GeeksforGeeks, IndiaBIX) require active internet. Fail in rural areas or during network disruptions."),
    ("High Ads & Distractions", "Online portals host heavy third-party banner ads, causing page delays, high RAM consumption, and user distraction."),
    ("Lack of App Consolidation", "Browser-only platforms are prone to notification triggers, taking focus away from study schedules."),
    ("Insecure Practice Sandbox", "Traditional study webs permit copy-pasting code, viewing developer console inspector (F12), and changing tabs, leading to inaccurate assessment statistics.")
]
add_bullet_slide(prs, "Existing System Analysis", existing_bullets)

# Slide 5: Proposed System
proposed_bullets = [
    ("Offline Integration", "Direct asset packaging under Android assets folder. Zero loading latencies, 100% offline-ready operations."),
    ("Consolidated Shell", "Jetpack Compose hosts a secure, hardware-accelerated WebView engine with custom permissions."),
    ("Secure Biometric Shield", "Protects user data files using Android biometrics (Fingerprint/PIN prompts) prior to main app loading."),
    ("Cheating Protection Suite", "Monitors keyboard actions (blocking copy-paste, print-screen, inspector shortcut commands) and auto-submits mock exams immediately when browser blur or tab swap is detected.")
]
add_bullet_slide(prs, "Proposed System Solution", proposed_bullets)

# Slide 6: Feasibility Study
feas_technical = [
    ("Technical Viability", "Built on top of stable native WebView frameworks and standard HTML5/CSS3/JS execution engines."),
    ("Integration Stability", "Native Compose security wrapper works seamlessly with modern Android API layers.")
]
feas_operational = [
    ("Economic Viability", "Serverless architecture leads to $0 deployment overhead. Uses open-source compilers and SDK resources."),
    ("Operational Viability", "Extremely intuitive UI layouts. Self-explanatory modules and automated timers ensure straightforward adoption.")
]
add_two_column_slide(prs, "Feasibility Study", "Technical & System Feasibility", feas_technical, "Economic & Operational Feasibility", feas_operational)

# Slide 7: Modules Description
modules_bullets = [
    ("Quantitative Aptitude", "Features 14 topics (Algebra, Speed, Work, Percentages, etc.). Concepts combined with 1000+ randomized questions."),
    ("Logical Reasoning", "Includes 9 modules (Blood Relations, Coding-Decoding, Seating Arrangements, etc.) highlighting reasoning rules."),
    ("Verbal Ability", "Covers 7 topics (Error Spotting, Para Jumbles, Grammar, Reading Comprehension) for vocabulary and grammar proficiency."),
    ("Coding Workspace", "Houses a JavaScript sandbox wrapper that processes and validates candidate scripts on the fly."),
    ("Mock Test Simulator", "Configurable exam generator (Easy, Medium, Hard, Full-Length) executing cheat protections and calculating user ratings.")
]
add_bullet_slide(prs, "Project Modules Description", modules_bullets)

# Slide 8: System Architecture
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide8, use_watermark=True)
add_slide_title(slide8, "System Architecture Diagram")

# Drawing vector diagram cards
def add_card(slide, title, text, left, top, width, height, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.1)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    
    if text:
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_SILVER
        p2.font.name = "Segoe UI"
        p2.alignment = PP_ALIGN.CENTER
    return shape

# User Layer (Left)
add_card(slide8, "User Interface Client", "Mobile Screen\nTouch Interactions\nKeyboard Input", 0.8, 2.5, 2.8, 3.2, ACCENT_PURPLE)

# Middle Container Stack (Android Shell & WebView Host)
add_card(slide8, "Android Native App Container", "Kotlin Activity Shell\nJetpack Compose States\nBiometrics SDK Lock Prompt", 4.8, 1.8, 3.8, 2.0, ACCENT_GREEN)
add_card(slide8, "WebView Web Host Engine", "DOM Caches Enabled\nFile Access Permitted\nJS Execution Sandbox", 4.8, 4.4, 3.8, 2.0, ACCENT_GREEN)

# Offline Database Layer (Right)
add_card(slide8, "LocalStorage State DB", "placement_prep_state\nBookmarks Index\nWrong Answers Log\nSolved Question Logs", 9.8, 2.5, 2.8, 3.2, ACCENT_PURPLE)

# Connectors/Arrows
# Arrow 1: User to Native app
a1 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.8), Inches(2.7), Inches(0.8), Inches(0.3))
a1.fill.solid()
a1.fill.fore_color.rgb = ACCENT_PURPLE
a1.line.fill.background()

# Arrow 2: User to Web wrapper directly (alternate path)
a2 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.8), Inches(5.1), Inches(0.8), Inches(0.3))
a2.fill.solid()
a2.fill.fore_color.rgb = ACCENT_PURPLE
a2.line.fill.background()

# Arrow 3: Native Container down to WebView Host
a3 = slide8.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.6), Inches(3.9), Inches(0.3), Inches(0.4))
a3.fill.solid()
a3.fill.fore_color.rgb = ACCENT_GREEN
a3.line.fill.background()

# Arrow 4: WebView Host to Storage DB
a4 = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.8), Inches(5.1), Inches(0.8), Inches(0.3))
a4.fill.solid()
a4.fill.fore_color.rgb = ACCENT_PURPLE
a4.line.fill.background()

# Slide 9: Dataflow Diagram (DFD)
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide9, use_watermark=True)
add_slide_title(slide9, "Level 1 Dataflow Diagram (DFD)")

def add_process_node(slide, number, name, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_BG
    shape.line.color.rgb = ACCENT_GREEN
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{number}\n{name}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return shape

# DFD Layout: Linear Flow with Label annotations
# Entity
user_ent = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.8), Inches(1.5), Inches(1.5))
user_ent.fill.solid()
user_ent.fill.fore_color.rgb = PANEL_BG
user_ent.line.color.rgb = ACCENT_PURPLE
user_ent.line.width = Pt(2)
user_ent.text_frame.paragraphs[0].text = "Candidate"
user_ent.text_frame.paragraphs[0].font.size = Pt(14)
user_ent.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE

# Processes
add_process_node(slide9, "1.0", "Validate\nBiometrics", 2.8, 2.7, 1.6, 1.6)
add_process_node(slide9, "2.0", "Execute\nPractice Board", 5.2, 2.7, 1.6, 1.6)
add_process_node(slide9, "3.0", "Monitor\nCheating Checks", 7.6, 2.7, 1.6, 1.6)
add_process_node(slide9, "4.0", "Synchronize\nState Logs", 10.0, 2.7, 1.6, 1.6)

# Data Store (Parallel lines / Open rect style box)
ds = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.0), Inches(2.95), Inches(1.1), Inches(1.2))
ds.fill.solid()
ds.fill.fore_color.rgb = BG_COLOR
ds.line.color.rgb = ACCENT_PURPLE
ds.line.width = Pt(1.5)
tf_ds = ds.text_frame
p_ds = tf_ds.paragraphs[0]
p_ds.text = "D1: Local\nStorage"
p_ds.font.size = Pt(12)
p_ds.font.bold = True
p_ds.font.color.rgb = TEXT_WHITE
p_ds.alignment = PP_ALIGN.CENTER

# Arrows and Annotations
arrow_x_coords = [2.2, 4.5, 6.9, 9.3, 11.7]
arrow_labels = ["Launch Request", "Success Verification", "Select Exam Mode", "Exam Results Log", "Update Schema"]

for idx, x in enumerate(arrow_x_coords):
    arr = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.35), Inches(0.5), Inches(0.2))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT_GREEN
    arr.line.fill.background()
    
    # Add label
    lbl = slide9.shapes.add_textbox(Inches(x - 0.4), Inches(2.1), Inches(1.3), Inches(0.5))
    lbl.text_frame.word_wrap = True
    p_lbl = lbl.text_frame.paragraphs[0]
    p_lbl.text = arrow_labels[idx]
    p_lbl.font.size = Pt(9)
    p_lbl.font.color.rgb = TEXT_SILVER
    p_lbl.alignment = PP_ALIGN.CENTER

# Slide 10: User Interface Design
ui_bullets = [
    ("Glassmorphism Concept", "Employs high-contrast translucent cards (`glass-card` styling) with fine borders and dark overlay shadows to present a polished, premium aesthetic."),
    ("Ambient Glow Fields", "Incorporates subtle background gradients and radial glows (`ambient-glow-1/2`) to create layout depth, keeping page elements readable."),
    ("Loading Screen Animation", "Combines custom CSS keyframes (`loading-screen`) with rotational spinner states and quotes, preventing visual flickering during initial app lock launch."),
    ("Mobile Side Drawer Menu", "Implements layout adapters and overlays, enabling standard navigation behaviors on smaller tablet and smartphone screens.")
]
add_bullet_slide(prs, "User Interface Design (SYSTEM DESIGN)", ui_bullets)

# Slide 11: Software Modules
sw_mod_bullets = [
    ("Kotlin Shell Module", "Hosts AppContent controllers, managing transitions between AppState.Splash, AppState.Lock, and AppState.Main WebView configurations."),
    ("Centralized State (storage.js)", "Consolidates getState(), updateState(), addBookmark(), and logSolvedQuestion() state hooks into a single global window API."),
    ("Exam Simulator Engine (mocktest.js)", "Constructs aggregated topic pools on demand. Leverages MockDataGen constructors to slice question arrays according to test tiers."),
    ("Security Shield (cheating_protection.js)", "Intercepts and restructures document inputs, suppressing ContextMenu, Copy/Cut/Paste hooks, and logging visibility switch events.")
]
add_bullet_slide(prs, "Software Modules Division", sw_mod_bullets)

# Slide 12: Database Design
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide12, use_watermark=True)
add_slide_title(slide12, "Database Schema Design (LocalStorage JSON)")

# Add table to represent localStorage fields
rows, cols = 5, 3
left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.5)
table_shape = slide12.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.3)
table.columns[2].width = Inches(7.23)

headers = ["Key Field Name", "Storage Data Type", "Detailed Description & Schema Role"]
schema_data = [
    ["theme", "String ('dark' | 'light')", "Saves candidate's preferred system interface color scheme settings (defaulting to dark)."],
    ["bookmarks", "Array of JSON Objects", "Tracks questions pinned by user. Indexes standard parameters (questionId, prompt, options, answer)."],
    ["wrongAnswers", "Array of Strings", "Logs IDs of questions answered incorrectly, populating the 'Retry Hub' queue for review sessions."],
    ["history", "Array of Solved Logs", "Stores log objects containing: { questionId, topic, subject, correct: Boolean, timeSpent: Int, timestamp }."]
]

# Style and populate headers
for c in range(3):
    cell = table.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_PURPLE
    p = cell.text_frame.paragraphs[0]
    p.text = headers[c]
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

# Populate rows
for r in range(1, 5):
    for c in range(3):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_BG
        p = cell.text_frame.paragraphs[0]
        p.text = schema_data[r-1][c]
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE if c < 2 else TEXT_SILVER
        p.alignment = PP_ALIGN.CENTER if c < 2 else PP_ALIGN.LEFT

# Slide 13: System Implementation
impl_bullets = [
    ("Asset Packaging Structure", "Assets copied to Android's `src/main/assets/appti` folder. Local WebView instances load filepaths offline (e.g. `file:///android_asset/appti/index.html`)."),
    ("Compose Interop Layer", "AndroidView compose block constructs layout containers. Attaches BackHandler checks to prevent unintentional app closing during back gestures."),
    ("Hardware Acceleration", "Enables hardware-accelerated WebView properties (allowing universal access flags and offline storage APIs) within the secure container.")
]
add_bullet_slide(prs, "System Implementation Overview", impl_bullets)

# Slide 14: Coding (Original Source Code)
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide14, use_watermark=True)
add_slide_title(slide14, "Original Code Implementation Snippets")

def add_code_panel(slide, code_text, left, top, width, height, title):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(28, 27, 34)
    bg.line.color.rgb = ACCENT_PURPLE
    bg.line.width = Pt(1.5)
    
    # Title
    header = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    p_h = header.text_frame.paragraphs[0]
    p_h.text = "  " + title
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = ACCENT_GREEN
    p_h.font.name = "Consolas"
    
    # Code Content
    content = slide.shapes.add_textbox(Inches(left), Inches(top + 0.35), Inches(width), Inches(height - 0.35))
    content.text_frame.word_wrap = True
    p_c = content.text_frame.paragraphs[0]
    p_c.text = code_text
    p_c.font.size = Pt(9.5)
    p_c.font.color.rgb = TEXT_WHITE
    p_c.font.name = "Consolas"

kotlin_snippet = """// MainActivity.kt - Screen Capture & Biometric Lock Setup
override fun onCreate(savedInstanceState: Bundle?) {
    super.onCreate(savedInstanceState)
    window.setFlags(
        WindowManager.LayoutParams.FLAG_SECURE,
        WindowManager.LayoutParams.FLAG_SECURE
    )
    enableEdgeToEdge()
    setContent {
        VetriPathTheme {
            Surface(modifier = Modifier.fillMaxSize()) { AppContent() }
        }
    }
}"""

js_snippet = """// cheating_protection.js - Browser Tab Blur Auto-Submission
function triggerFocusLossWarning() {
    if (!active) return;
    
    PlacementPrepState.dispatchToast(
        "🚨 Focus loss/Tab switch detected! Exam auto-submitted.", 
        "danger"
    );
    disable();
    if (limitReachedCallback) {
        limitReachedCallback(); // Auto-submit hook
    }
}"""

add_code_panel(slide14, kotlin_snippet, 0.7, 1.8, 5.7, 5.0, "Kotlin (Android Native Screen Protection)")
add_code_panel(slide14, js_snippet, 6.9, 1.8, 5.7, 5.0, "JavaScript (Web Cheating Controller)")

# Slide 15: Screenshots Mockups
mock_bullets = [
    ("Dashboard Screens", "Four statistics widgets displaying Chapter quantities, total Questions (2000+), adaptable levels, and 100% offline operational flags."),
    ("Biometric Lock Prompt", "Native OS overlay displaying Fingerprint sensor icon, requesting authentication with cancel options."),
    ("Practice Board Console", "Split column showing study/shortcuts on the left, and multi-choice question interface with bookmark toggles on the right."),
    ("Mock Exam Workspace", "Countdown timer bar at the top, navigation board tracking answered questions, and immediate submission handlers.")
]
add_bullet_slide(prs, "System Interface Visualizations", mock_bullets)

# Slide 16: System Testing
test_col1 = [
    ("Biometric Locking Test", "Ensures Splash transitions to Lockscreen. Verifies successful scanner validation, PIN fallbacks, and sensor-missing bypass cases."),
    ("Offline Operations Test", "Confirms WebView loads correctly when mobile devices are on airplane mode, verifying storage sync works offline.")
]
test_col2 = [
    ("Cheating Prevention Test", "Simulates window blur and visibility toggles during active mock tests. Validates instant test submission triggers."),
    ("UI Responsive Test", "Checks HTML grid rendering, fonts scaling, and menu navigation on viewport sizes from 5\" devices to 10\" tablet screens.")
]
add_two_column_slide(prs, "System Verification & Quality Testing", "Biometric & Offline Testing Logs", test_col1, "Exam Integrity & Responsive Testing Logs", test_col2)

# Slide 17: Future Enhancements
future_bullets = [
    ("Cloud Sync Middleware", "Integrates standard Google Drive REST APIs for multi-device sync, backup, and restore."),
    ("AI Personal Analytics", "Recommends prep paths based on solved question difficulty scores and weaker subjects."),
    ("Peer-to-Peer Competitive Arena", "Online testing rooms and real-time scoreboards using standard WebSockets."),
    ("Instructor Management Console", "Allows educators to build assessment profiles, export progress reports, and track candidate metrics.")
]
add_bullet_slide(prs, "Future System Enhancements", future_bullets)

# Slide 18: Summary & Key Accomplishments
summary_bullets = [
    ("Self-Contained Hybrid Portal", "Successfully developed a hybrid application framework that runs offline without loading delays."),
    ("High-Grade UI/UX Design", "Custom CSS and glassmorphism elements matching standard premium web layouts."),
    ("Robust Device Security", "Implemented biometrics locking to keep user logs protected locally on the device."),
    ("Assessment Integrity Control", "Built automated anti-cheating systems, ensuring test evaluations are reliable.")
]
add_bullet_slide(prs, "Project Summary & Key Accomplishments", summary_bullets)

# Slide 19: Thank You Slide
slide19 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide19, use_watermark=False)
if os.path.exists(logo_path):
    slide19.shapes.add_picture(logo_path, Inches(5.41), Inches(1.2), width=Inches(2.5), height=Inches(2.5))

tx19 = slide19.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(3.0))
tf19 = tx19.text_frame
tf19.word_wrap = True

p_ty = tf19.paragraphs[0]
p_ty.text = "THANK YOU"
p_ty.font.size = Pt(54)
p_ty.font.bold = True
p_ty.font.color.rgb = TEXT_WHITE
p_ty.font.name = "Segoe UI"
p_ty.alignment = PP_ALIGN.CENTER
p_ty.space_after = Pt(12)

p_tag = tf19.add_paragraph()
p_tag.text = '"Walk the Path. Become the Victory. Inspire the World."'
p_tag.font.size = Pt(20)
p_tag.font.italic = True
p_tag.font.color.rgb = ACCENT_GREEN
p_tag.font.name = "Segoe UI"
p_tag.alignment = PP_ALIGN.CENTER
p_tag.space_after = Pt(24)

p_info = tf19.add_paragraph()
p_info.text = "VetriPath Team • App & Website Project"
p_info.font.size = Pt(14)
p_info.font.color.rgb = TEXT_SILVER
p_info.font.name = "Segoe UI"
p_info.alignment = PP_ALIGN.CENTER

# Save final presentation
output_name = "VetriPath_Project_Presentation.pptx"
prs.save(output_name)
print(f"Presentation saved successfully as {output_name}.")

# Cleanup faded logo path
if os.path.exists(faded_logo_path):
    os.remove(faded_logo_path)
    print("Cleaned up temporary files.")
