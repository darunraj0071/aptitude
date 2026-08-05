import os
import sys
import qrcode
from PIL import Image, ImageDraw, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Paths Setup
base_dir = r"c:\Users\sanja\Downloads\appti (2) (1)\appti (3)"
mca_img_dir = os.path.join(base_dir, "appti", "images", "mca")
screenshots_dir = os.path.join(base_dir, "screenshots")
logo_path = os.path.join(base_dir, "appti", "images", "logo.png")
faded_logo_path = os.path.join(base_dir, "faded_logo.png")
qr_path = os.path.join(base_dir, "vetripath_qr.png")

output_ppt = os.path.join(base_dir, "VetriPath_Project_Presentation.pptx")
output_ppt_alt = os.path.join(base_dir, "VetriPath_Learn_MCA_Mini_Project_Presentation.pptx")

# Generate High-Error-Correction Seamlessly Merged Circular Logo QR Code
try:
    qr = qrcode.QRCode(
        version=4,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=12,
        border=2
    )
    qr.add_data("https://www.vetripathlearn.site/")
    qr.make(fit=True)
    qr_img = qr.make_image(fill_color="#00C8FF", back_color="#0B0B0B").convert("RGBA")
    
    if os.path.exists(logo_path):
        logo = Image.open(logo_path).convert("RGBA")
        qr_w, qr_h = qr_img.size
        l_size = int(qr_w // 3.5)
        
        mask = Image.new("L", (l_size, l_size), 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, l_size, l_size), fill=255)
        
        logo_resized = logo.resize((l_size, l_size), Image.Resampling.LANCZOS)
        
        circular_bg = Image.new("RGBA", (l_size + 12, l_size + 12), (0, 0, 0, 0))
        draw_bg = ImageDraw.Draw(circular_bg)
        draw_bg.ellipse((0, 0, l_size + 12, l_size + 12), fill="#0B0B0B", outline="#00C8FF", width=3)
        
        pos = ((qr_w - l_size) // 2, (qr_h - l_size) // 2)
        pos_bg = (pos[0] - 6, pos[1] - 6)
        
        qr_img.paste(circular_bg, pos_bg, circular_bg)
        qr_img.paste(logo_resized, pos, mask)
        
    qr_img.save(qr_path)
    print("Seamlessly Merged Circular Logo QR Code created successfully.")
except Exception as e:
    print(f"Branded QR Code creation error: {e}")

# Generate Faded Watermark Logo Image
try:
    if os.path.exists(logo_path):
        logo_img = Image.open(logo_path).convert("RGBA")
        r, g, b, alpha = logo_img.split()
        alpha_faded = ImageEnhance.Brightness(alpha).enhance(0.12)
        faded_logo = Image.merge("RGBA", (r, g, b, alpha_faded))
        faded_logo.save(faded_logo_path)
except Exception as e:
    print(f"Watermark creation error: {e}")

# Initialize Widescreen Presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Color Palette
BG_COLOR = RGBColor(11, 11, 11)       # #0B0B0B Dark Matte Black
CARD_BG = RGBColor(21, 21, 26)       # #15151A Translucent Glass Container
CARD_BORDER = RGBColor(0, 200, 255)  # #00C8FF Electric Cyan Border
ACCENT_CYAN = RGBColor(0, 229, 255)  # #00E5FF Neon Cyan Accent
ACCENT_INDIGO = RGBColor(79, 70, 229)# #4F46E5 Indigo Accent
TEXT_WHITE = RGBColor(255, 255, 255) # White Text
TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8 Silver Muted Text
HEADER_CYAN = RGBColor(56, 189, 248) # #38BDF8 Header Cyan

def add_slide_background_and_transition(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # 1. Slide Fade Transition (Duration: 0.5s, On Mouse Click only)
    try:
        trans_xml = parse_xml('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>')
        slide.element.append(trans_xml)
    except Exception as e:
        pass
    
    # 2. Background Centered Faded Logo Watermark
    if os.path.exists(faded_logo_path):
        slide.shapes.add_picture(faded_logo_path, Inches(4.666), Inches(1.75), width=Inches(4.0), height=Inches(4.0))
        
    # Header glowing top border line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = ACCENT_CYAN
    top_line.line.color.rgb = ACCENT_CYAN
    
    # Footer text and divider
    footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(7.0), Inches(11.933), Inches(0.02))
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = RGBColor(51, 65, 85)
    footer_line.line.color.rgb = RGBColor(51, 65, 85)
    
    f_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(11.933), Inches(0.35))
    tf_f = f_box.text_frame
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
    p_f = tf_f.paragraphs[0]
    p_f.text = "VetriPath Learn — Master of Computer Applications (MCA) Mini Project Presentation"
    p_f.font.size = Pt(10)
    p_f.font.color.rgb = TEXT_MUTED
    p_f.font.name = "Segoe UI"

def add_header(slide, title_text, category_text="VETRIPATH LEARN MCA PRESENTATION"):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(12.0), Inches(0.25), width=Inches(0.65), height=Inches(0.65))
        
    cat_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(11.0), Inches(0.3))
    tf_c = cat_box.text_frame
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = category_text.upper()
    p_c.font.size = Pt(10.5)
    p_c.font.bold = True
    p_c.font.color.rgb = HEADER_CYAN
    p_c.font.name = "Segoe UI"
    
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.0), Inches(0.6))
    tf_t = title_box.text_frame
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(24)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.font.name = "Segoe UI"
    return title_box.shape_id

def create_glass_card(slide, left, top, width, height, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

def add_sequential_click_animations(slide, shape_ids):
    """
    Appends professional PowerPoint OpenXML timing sequence so every element
    appears one-by-one upon user mouse click (On Click, 0.5s duration, 0s delay).
    """
    if not shape_ids:
        return
        
    child_nodes = ""
    ctn_id = 10
    for idx, spid in enumerate(shape_ids):
        child_nodes += f'''
        <p:par>
          <p:cTn id="{ctn_id}" fill="hold">
            <p:stCondLst>
              <p:cond delay="0"/>
            </p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{ctn_id+1}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                  <p:stCondLst>
                    <p:cond delay="0"/>
                  </p:stCondLst>
                  <p:childTnLst>
                    <p:set>
                      <p:cb>
                        <p:spTarget spid="{spid}"/>
                      </p:cb>
                      <p:to>
                        <p:visibility val="visible"/>
                      </p:to>
                    </p:set>
                    <p:anim filter="blinds(horizontal)" calcmode="lin" valueType="num">
                      <p:cb>
                        <p:spTarget spid="{spid}"/>
                      </p:cb>
                    </p:anim>
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>
        '''
        ctn_id += 5
        
    timing_xml_str = f'''
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="root">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    {child_nodes}
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    '''
    try:
        slide.element.append(parse_xml(timing_xml_str))
    except Exception as e:
        print(f"Timing animation warning: {e}")

print("Building Presentation Slides with Professional On-Click Sequential Animations & Fade Transitions...")

# SLIDE 1: TITLE SLIDE
slide1 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide1)

hero_card = create_glass_card(slide1, Inches(0.7), Inches(0.65), Inches(11.933), Inches(6.1), ACCENT_CYAN)

if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(1.1), Inches(0.95), width=Inches(1.2), height=Inches(1.2))

t_box = slide1.shapes.add_textbox(Inches(2.5), Inches(0.9), Inches(9.8), Inches(2.5))
tf = t_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "VetriPath Learn"
p1.font.size = Pt(42)
p1.font.bold = True
p1.font.color.rgb = ACCENT_CYAN
p1.font.name = "Segoe UI"

p2 = tf.add_paragraph()
p2.text = "Offline-First Hybrid Career Guidance, Aptitude Learning and Placement Preparation Platform"
p2.font.size = Pt(19)
p2.font.bold = True
p2.font.color.rgb = TEXT_WHITE
p2.font.name = "Segoe UI"
p2.space_before = Pt(6)

p3 = tf.add_paragraph()
p3.text = "A Mini Project Presentation Submitted for the Award of the Degree of Master of Computer Applications (MCA)"
p3.font.size = Pt(13)
p3.font.color.rgb = TEXT_MUTED
p3.font.name = "Segoe UI"
p3.space_before = Pt(6)

card_sub = create_glass_card(slide1, Inches(1.1), Inches(3.7), Inches(5.5), Inches(2.7), ACCENT_INDIGO)
box_sub = slide1.shapes.add_textbox(Inches(1.3), Inches(3.85), Inches(5.1), Inches(2.4))
tf_sub = box_sub.text_frame
tf_sub.word_wrap = True

ps0 = tf_sub.paragraphs[0]
ps0.text = "SUBMITTED BY:"
ps0.font.size = Pt(13.5)
ps0.font.bold = True
ps0.font.color.rgb = HEADER_CYAN

ps1 = tf_sub.add_paragraph()
ps1.text = "1. DARUN RAJ D (RA2532241040009)"
ps1.font.size = Pt(13)
ps1.font.bold = True
ps1.font.color.rgb = TEXT_WHITE
ps1.space_before = Pt(8)

ps2 = tf_sub.add_paragraph()
ps2.text = "2. NISMETHA P (RA2532241040021)"
ps2.font.size = Pt(13)
ps2.font.bold = True
ps2.font.color.rgb = TEXT_WHITE
ps2.space_before = Pt(6)

ps3 = tf_sub.add_paragraph()
ps3.text = "Department of Computer Science & Applications"
ps3.font.size = Pt(11.5)
ps3.font.color.rgb = TEXT_MUTED
ps3.space_before = Pt(6)

card_guide = create_glass_card(slide1, Inches(6.833), Inches(3.7), Inches(5.5), Inches(2.7), ACCENT_CYAN)
box_guide = slide1.shapes.add_textbox(Inches(7.033), Inches(3.85), Inches(5.1), Inches(2.4))
tf_guide = box_guide.text_frame
tf_guide.word_wrap = True

pg0 = tf_guide.paragraphs[0]
pg0.text = "GUIDED BY:"
pg0.font.size = Pt(13.5)
pg0.font.bold = True
pg0.font.color.rgb = HEADER_CYAN

pg1 = tf_guide.add_paragraph()
pg1.text = "DR. A. MEENAKSHI, M.C.A., M.Phil., M.S., Ph.D."
pg1.font.size = Pt(13)
pg1.font.bold = True
pg1.font.color.rgb = TEXT_WHITE
pg1.space_before = Pt(8)

pg2 = tf_guide.add_paragraph()
pg2.text = "Assistant Professor"
pg2.font.size = Pt(13)
pg2.font.bold = True
pg2.font.color.rgb = ACCENT_CYAN
pg2.space_before = Pt(4)

pg3 = tf_guide.add_paragraph()
pg3.text = "Department of Computer Science & Applications"
pg3.font.size = Pt(11.5)
pg3.font.color.rgb = TEXT_MUTED
pg3.space_before = Pt(6)

add_sequential_click_animations(slide1, [t_box.shape_id, card_sub.shape_id, card_guide.shape_id])

# SLIDE 2: PROJECT ABSTRACT
slide2 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide2)
t2_id = add_header(slide2, "Project Executive Abstract", "ABSTRACT")

card_abs = create_glass_card(slide2, Inches(0.7), Inches(1.3), Inches(11.933), Inches(5.4), ACCENT_CYAN)
tb_abs = slide2.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
tf_abs = tb_abs.text_frame
tf_abs.word_wrap = True

pa0 = tf_abs.paragraphs[0]
pa0.text = "VetriPath Learn — Offline-First Executive Project Abstract"
pa0.font.size = Pt(21)
pa0.font.bold = True
pa0.font.color.rgb = HEADER_CYAN

abs_bullets = [
    ("Industry Context & Core Issue", "In contemporary campus recruitment drives, candidates face severe hurdles due to network volatility, ad-distracted portals, server latency, and lack of assessment anti-cheating controls."),
    ("Dual-Platform Hybrid Architecture", "VetriPath Learn delivers a 100% offline-first dual-platform solution: a high-performance web portal for desktop practice, and a native Kotlin Jetpack Compose Android app (`VetriPath.apk`) for secure mobile practice."),
    ("Comprehensive Practice Workstation", "Houses 12 integrated modules providing 2,000+ quantitative aptitude, reasoning, verbal, interactive JavaScript code execution sandbox, and 6-month placement timelines."),
    ("Zero Server Reliance & $0 Cost", "Centralizes offline data persistence inside HTML5 LocalStorage (`storage.js`), operating self-sufficiently under Airplane Mode with $0 operational server costs.")
]

for title, desc in abs_bullets:
    p = tf_abs.add_paragraph()
    p.text = f"•  {title}:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_before = Pt(10)
    
    p2 = tf_abs.add_paragraph()
    p2.text = f"   {desc}"
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(2)

add_sequential_click_animations(slide2, [t2_id, card_abs.shape_id])

# SLIDE 3: AGENDA
slide3 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide3)
t3_id = add_header(slide3, "Presentation Agenda", "OVERVIEW")

agenda_items = [
    ("01", "Project Executive Abstract"),
    ("02", "Introduction & Project Background"),
    ("03", "Problem Statement & Industry Gaps"),
    ("04", "Our Solution & Technical Approach"),
    ("05", "What Problems Are Solved by VetriPath"),
    ("06", "Project Objectives & Key Goals"),
    ("07", "Existing System vs VetriPath Comparison"),
    ("08", "Drawbacks of Online Preparation Portals"),
    ("09", "Technology Stack Breakdown"),
    ("10", "Dual-Platform System Architecture"),
    ("11", "Core Software Modules & Descriptions"),
    ("12", "Database Design & LocalStorage Schema"),
    ("13", "UML Diagrams (DFD, ER, Use Case)"),
    ("14", "System Implementation & Native Shell"),
    ("15", "Application UI Screenshots Gallery"),
    ("16", "System Testing & QA Execution Matrix"),
    ("17", "Key System Advantages & Benefits"),
    ("18", "Future Roadmap & Project Synthesis")
]

s3_shapes = [t3_id]
for idx, (num, title) in enumerate(agenda_items):
    col = idx // 9
    row = idx % 9
    x = Inches(0.7 + col * 6.0)
    y = Inches(1.25 + row * 0.62)
    card = create_glass_card(slide3, x, y, Inches(5.7), Inches(0.56), CARD_BORDER)
    s3_shapes.append(card.shape_id)
    
    tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.06), Inches(5.3), Inches(0.42))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}.  {title}"
    p.font.size = Pt(12.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"

add_sequential_click_animations(slide3, s3_shapes)

# SLIDE 4: INTRODUCTION
slide4 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide4)
t4_id = add_header(slide4, "Introduction & Project Background", "PROJECT CONTEXT")

intro_cards = [
    ("Placement Environment", "Corporate placement recruitment drives evaluate candidates across quantitative aptitude, logical reasoning, verbal ability, algorithmic coding, and interview screening tests.", "• Multi-stage campus selection criteria\n• High speed & accuracy expectations\n• Requires continuous daily practice"),
    ("Rural Network Bottleneck", "Students face severe network drops. Online portals rely on persistent cloud servers and heavy ad banners, causing page loading failures and loss of practice records.", "• Severe internet drops in tier-2/3 towns\n• Loss of test submissions & answers\n• Distractive ad banner clutter"),
    ("VetriPath Dual-Platform", "Delivers a 100% offline-first solution: a high-performance web portal for desktop practice, and a native Kotlin Jetpack Compose Android app (`VetriPath.apk`) for secure mobile practice.", "• 100% Airplane Mode operational\n• Zero cloud server dependencies\n• Sub-50ms instant local rendering")
]

s4_shapes = [t4_id]
for idx, (title, desc, bullet) in enumerate(intro_cards):
    x = Inches(0.7 + idx * 3.95)
    card = create_glass_card(slide4, x, Inches(1.3), Inches(3.75), Inches(5.4), ACCENT_CYAN)
    s4_shapes.append(card.shape_id)
    
    tb = slide4.shapes.add_textbox(x + Inches(0.25), Inches(1.5), Inches(3.25), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = bullet
    p3.font.size = Pt(11.5)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(12)

add_sequential_click_animations(slide4, s4_shapes)

# SLIDE 5: PROBLEM STATEMENT
slide5 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide5)
t5_id = add_header(slide5, "Problem Statement", "INDUSTRY CHALLENGES")

problems = [
    ("01", "Persistent Network Dependency", "Online portals fail during internet drops, causing answer submission errors and progress data loss.", "Impact: Candidates lose exam state & mock test scores."),
    ("02", "Fragmented Preparation Workstations", "Candidates must navigate separate websites for aptitude, separate compilers for coding, and endure ad clutter.", "Impact: Inefficient preparation workflow & lost time."),
    ("03", "Lack of Mobile Security Control", "Standard browsers lack native device locking, permitting screenshot captures and background app switching.", "Impact: Unauthenticated practice & security risks."),
    ("04", "Unprotected Assessment Integrity", "Mock tests permit copy-pasting, tab switching, and window minimization without issuing security alerts.", "Impact: Unfair assessment results & exam cheating.")
]

s5_shapes = [t5_id]
for idx, (num, title, desc, impact) in enumerate(problems):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide5, x, y, Inches(5.75), Inches(2.55), ACCENT_INDIGO)
    s5_shapes.append(card.shape_id)
    
    tb = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(0.18), Inches(5.15), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"{num}. {title}"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(6)
    
    p3 = tf.add_paragraph()
    p3.text = impact
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = HEADER_CYAN
    p3.space_before = Pt(6)

add_sequential_click_animations(slide5, s5_shapes)

# SLIDE 6: OUR SOLUTION
slide6 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide6)
t6_id = add_header(slide6, "Our Solution — VetriPath Learn", "PROPOSED SOLUTION")

sol_cards = [
    ("100% Offline-First Web PWA", "High-performance HTML5/CSS3/ES6 web application operating completely offline with local JavaScript engines."),
    ("Hardened Native Android APK", "Kotlin 1.9 & Jetpack Compose WebView shell (`VetriPath.apk`) enforcing native biometric security & window protection."),
    ("Unified All-in-One Workstation", "Houses 2,000+ quantitative aptitude, reasoning, verbal, interactive coding runner, and placement roadmaps in one portal."),
    ("Serverless $0 State Engine", "Centralizes offline progress management inside LocalStorage (`storage.js`), eliminating cloud database fees.")
]

s6_shapes = [t6_id]
for idx, (title, desc) in enumerate(sol_cards):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide6, x, y, Inches(5.75), Inches(2.55), ACCENT_CYAN)
    s6_shapes.append(card.shape_id)
    
    tb = slide6.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.15), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"💡 {title}"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)

add_sequential_click_animations(slide6, s6_shapes)

# SLIDE 7: WHAT ARE THE PROBLEMS SOLVED
slide7 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide7)
t7_id = add_header(slide7, "What Problems Are Solved by VetriPath", "PROBLEM SOLVED")

solved_items = [
    ("Network Drop & Data Loss Solved", "PROBLEM: Portals crash on internet drop losing exam state.\nSOLVED: 100% Airplane Mode execution guarantees zero answer loss or submission failures."),
    ("Distractive Ad Clutter Solved", "PROBLEM: Portals load heavy ad banners causing 5-sec page delays.\nSOLVED: Sub-50ms instant local DOM rendering with zero third-party advertisements."),
    ("Exam Copy-Pasting & Cheating Solved", "PROBLEM: Candidates switch tabs to search answers online during exams.\nSOLVED: Focus loss guard (`cheating_protection.js`) auto-submits exam on tab blur."),
    ("Mobile Screen Capture & Leak Solved", "PROBLEM: Mobile browsers permit recording & screenshot captures.\nSOLVED: Native Android `FLAG_SECURE` completely blocks PrintScreen & screen recording.")
]

s7_shapes = [t7_id]
for idx, (title, desc) in enumerate(solved_items):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide7, x, y, Inches(5.75), Inches(2.55), CARD_BORDER)
    s7_shapes.append(card.shape_id)
    
    tb = slide7.shapes.add_textbox(x + Inches(0.3), y + Inches(0.18), Inches(5.15), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"✅ {title}"
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)

add_sequential_click_animations(slide7, s7_shapes)

# SLIDE 8: OBJECTIVES
slide8 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide8)
t8_id = add_header(slide8, "Project Objectives", "KEY GOALS")

objectives = [
    ("100% Offline-First Execution", "Engineered dual-platform architecture (Web & Android APK) operating seamlessly under complete Airplane Mode."),
    ("2,000+ Question Workstation", "Provide randomized quantitative aptitude, reasoning, verbal, and client-side JavaScript coding sandbox."),
    ("Client-Side Anti-Cheating Suite", "Build `cheating_protection.js` with window focus loss detection, auto-submission, and shortcut blocking."),
    ("Native Android Biometric Shield", "Implement Jetpack Compose BiometricPrompt scanner and WindowManager `FLAG_SECURE` screenshot prevention."),
    ("Serverless $0 Storage Architecture", "Centralize state management inside LocalStorage (`storage.js`), eliminating cloud server fees ($0 cost).")
]

s8_shapes = [t8_id]
for idx, (title, desc) in enumerate(objectives):
    y = Inches(1.25 + idx * 1.1)
    card = create_glass_card(slide8, Inches(0.7), y, Inches(11.933), Inches(0.98), CARD_BORDER)
    s8_shapes.append(card.shape_id)
    
    tb = slide8.shapes.add_textbox(Inches(1.0), y + Inches(0.12), Inches(11.3), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"🎯  {title}"
    p.font.size = Pt(15.5)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(3)

add_sequential_click_animations(slide8, s8_shapes)

# SLIDE 9: EXISTING SYSTEM COMPARATIVE STUDY
slide9 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide9)
t9_id = add_header(slide9, "Existing System Analysis", "COMPARATIVE STUDY")

comp_headers = ["Evaluation Metric", "Online Portals (IndiaBIX/GFG)", "VetriPath Hybrid Platform"]
comp_rows = [
    ["Network Dependency", "Requires active continuous internet connection", "100% Offline-First (Airplane Mode operational)"],
    ["Page Load Speed", "2,000ms - 5,000ms (Ad banner loading delays)", "Sub-50ms instant local WebView rendering"],
    ["Mobile Device Security", "Unauthenticated public browser history", "Native Jetpack Compose Biometric Sensor Lock"],
    ["Exam Cheating Monitor", "Permits copy-paste & tab switching", "Auto-submits exam on tab blur / focus loss"],
    ["Deployment Overhead", "High monthly cloud server & DB costs", "$0 serverless client-side storage architecture"]
]

table_shape = slide9.shapes.add_table(6, 3, Inches(0.7), Inches(1.3), Inches(11.933), Inches(5.4))
table = table_shape.table
table.columns[0].width = Inches(3.2)
table.columns[1].width = Inches(4.3)
table.columns[2].width = Inches(4.433)

for c_idx, h_text in enumerate(comp_headers):
    cell = table.cell(0, c_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_INDIGO
    p = cell.text_frame.paragraphs[0]
    p.text = h_text
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

for r_idx, row in enumerate(comp_rows):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if c_idx < 2 else RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_WHITE if c_idx == 2 else TEXT_MUTED

add_sequential_click_animations(slide9, [t9_id, table_shape.shape_id])

# SLIDE 10: DRAWBACKS OF EXISTING SYSTEM
slide10 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide10)
t10_id = add_header(slide10, "Drawbacks of Existing Systems", "SYSTEMIC LIMITATIONS")

drawbacks = [
    ("Network Vulnerability", "Portals fail to render or save user test submissions during internet connectivity drops.", "• Unsaved test data\n• Frequent page timeouts"),
    ("Intrusive Ad Clutter", "Excessive third-party ad banners degrade rendering speed and cause severe candidate distraction.", "• Slow DOM loading\n• Reduced visual focus"),
    ("Unprotected Security", "Standard mobile web browsers lack native biometric locking and permit screenshot captures.", "• Unrestricted app access\n• Easy screen recording"),
    ("Unprotected Assessments", "Mock tests permit copy-pasting, tab switching, and window minimization without issuing security alerts.", "• Exam copy-pasting\n• Lack of tab blur alerts")
]

s10_shapes = [t10_id]
for idx, (title, desc, detail) in enumerate(drawbacks):
    x = Inches(0.7 + idx * 2.98)
    card = create_glass_card(slide10, x, Inches(1.3), Inches(2.8), Inches(5.4), CARD_BORDER)
    s10_shapes.append(card.shape_id)
    
    tb = slide10.shapes.add_textbox(x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"⚠️\n{title}"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = detail
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(12)

add_sequential_click_animations(slide10, s10_shapes)

# SLIDE 11: TECHNOLOGY STACK BREAKDOWN
slide11 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide11)
t11_id = add_header(slide11, "Technology Stack Breakdown", "TECH ARCHITECTURE")

tech_stack = [
    ("Web Frontend Tier", "HTML5, CSS3 Glassmorphism, ES6 JavaScript, Tailwind CSS v3, `canvas-3d.js` particles"),
    ("Mobile Native Shell", "Kotlin 1.9, Jetpack Compose 1.5, Android WebView Interop (`setDomStorageEnabled`)"),
    ("Native Security SDK", "AndroidX Biometric SDK 1.2, WindowManager `FLAG_SECURE` Flags"),
    ("Offline Storage Engine", "HTML5 Web Storage API (`window.localStorage`), Centralized State Manager (`storage.js`)"),
    ("Deployment Outputs", "Vercel Static Web PWA + Native Android Binary APK (`VetriPath.apk`)"),
    ("IDE & Build Tools", "Android Studio Hedgehog, VS Code, Git, GitHub Code Repository")
]

s11_shapes = [t11_id]
for idx, (category, details) in enumerate(tech_stack):
    col = idx % 3
    row = idx // 3
    x = Inches(0.7 + col * 3.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide11, x, y, Inches(3.75), Inches(2.55), ACCENT_INDIGO)
    s11_shapes.append(card.shape_id)
    
    tb = slide11.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2), Inches(3.25), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = category
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = details
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)

add_sequential_click_animations(slide11, s11_shapes)

# SLIDE 12: SYSTEM ARCHITECTURE DIAGRAM
slide12 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide12)
t12_id = add_header(slide12, "System Architecture Diagram", "4-TIER ARCHITECTURE")

arch_pic = None
arch_path = os.path.join(mca_img_dir, "architecture_diagram.png")
if os.path.exists(arch_path):
    arch_pic = slide12.shapes.add_picture(arch_path, Inches(0.7), Inches(1.3), Inches(7.5), Inches(5.4))

card12 = create_glass_card(slide12, Inches(8.4), Inches(1.3), Inches(4.233), Inches(5.4), ACCENT_CYAN)
tb = slide12.shapes.add_textbox(Inches(8.65), Inches(1.5), Inches(3.733), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Architecture Callouts:"
p.font.size = Pt(16.5)
p.font.bold = True
p.font.color.rgb = HEADER_CYAN

callouts = [
    "1. Client UI Layer: Responsive HTML5/Tailwind CSS viewports.",
    "2. Native Shell Layer: Kotlin Compose activity managing Biometric lock.",
    "3. Execution & Security: JS code runner & focus monitor.",
    "4. Persistence: LocalStorage JSON state manager (`storage.js`)."
]

for c in callouts:
    p2 = tf.add_paragraph()
    p2.text = c
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(12)

s12_seq = [t12_id]
if arch_pic: s12_seq.append(arch_pic.shape_id)
s12_seq.append(card12.shape_id)
add_sequential_click_animations(slide12, s12_seq)

# SLIDE 13: CORE SOFTWARE MODULES & DESCRIPTIONS
slide13 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide13)
t13_id = add_header(slide13, "Core Software Modules & Module Descriptions", "MODULE SPECIFICATIONS")

modules_grid = [
    ("1. Dashboard (`index.html`)", "Displays progress meters, question metrics (2000+), solved count, & action cards."),
    ("2. Aptitude (`aptitude.html`)", "Covers 14 mathematical topics with formula guides & choice feedback (`quiz.js`)."),
    ("3. Reasoning (`reasoning.html`)", "Analytical deduction across 9 topics including Blood Relations & Syllogisms."),
    ("4. Verbal (`verbal.html`)", "Enhances English grammar, error spotting rules, & reading comprehension."),
    ("5. Coding Sandbox (`coding.html`)", "Client-side JS code runner (`coding_runner.js`) with stdout console capture."),
    ("6. Practice Hub (`practice.html`)", "Central directory listing all practice topics with category tabs & search."),
    ("7. Mock Test (`mocktest.html`)", "Timed corporate exam simulator (`mocktest.js`) with randomized question pools."),
    ("8. Anti-Cheating Suite", "Monitors window blur & tab switching, triggering exam auto-submission."),
    ("9. Native Biometric Shield", "Enforces fingerprint lock & FLAG_SECURE screenshot prevention on Android."),
    ("10. Bookmarks Hub", "Maintains user pinned questions array in LocalStorage for candidate revision."),
    ("11. Placement Roadmap", "6-month corporate prep timeline checklist (`roadmap_data.js`)."),
    ("12. Global Search Index", "Real-time keyword indexing across all 30+ practice topics (`search.js`).")
]

s13_shapes = [t13_id]
for idx, (title, desc) in enumerate(modules_grid):
    col = idx % 4
    row = idx // 4
    x = Inches(0.7 + col * 2.95)
    y = Inches(1.3 + row * 1.85)
    card = create_glass_card(slide13, x, y, Inches(2.8), Inches(1.72), CARD_BORDER)
    s13_shapes.append(card.shape_id)
    
    tb = slide13.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), Inches(2.5), Inches(1.45))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(4)

add_sequential_click_animations(slide13, s13_shapes)

# SLIDE 14: DATABASE DESIGN & ER DIAGRAM
slide14 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide14)
t14_id = add_header(slide14, "Database Design & LocalStorage Schema", "DATA MODELING")

er_pic = None
er_path = os.path.join(mca_img_dir, "er_diagram.png")
if os.path.exists(er_path):
    er_pic = slide14.shapes.add_picture(er_path, Inches(0.7), Inches(1.3), Inches(7.5), Inches(5.4))

card14 = create_glass_card(slide14, Inches(8.4), Inches(1.3), Inches(4.233), Inches(5.4), ACCENT_INDIGO)
tb = slide14.shapes.add_textbox(Inches(8.65), Inches(1.5), Inches(3.733), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "LocalStorage JSON Keys:"
p.font.size = Pt(16.5)
p.font.bold = True
p.font.color.rgb = HEADER_CYAN

db_tables = [
    "• placement_prep_state: Root container object",
    "• theme: User visual mode ('dark'/'light')",
    "• bookmarks: Saved question objects array",
    "• wrongAnswers: Missed question retry queue",
    "• history: Chronological mock test result logs"
]

for t in db_tables:
    p2 = tf.add_paragraph()
    p2.text = t
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)

s14_seq = [t14_id]
if er_pic: s14_seq.append(er_pic.shape_id)
s14_seq.append(card14.shape_id)
add_sequential_click_animations(slide14, s14_seq)

# SLIDE 15: UML DIAGRAMS OVERVIEW
slide15 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide15)
t15_id = add_header(slide15, "Unified Modeling Language (UML) Overview", "SYSTEM UML")

uc_pic = None
uc_path = os.path.join(mca_img_dir, "use_case_diagram.png")
if os.path.exists(uc_path):
    uc_pic = slide15.shapes.add_picture(uc_path, Inches(0.7), Inches(1.3), Inches(8.0), Inches(5.4))

card15 = create_glass_card(slide15, Inches(8.9), Inches(1.3), Inches(3.733), Inches(5.4), ACCENT_CYAN)
tb = slide15.shapes.add_textbox(Inches(9.1), Inches(1.5), Inches(3.333), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Actor Use Cases:"
p.font.size = Pt(16.5)
p.font.bold = True
p.font.color.rgb = HEADER_CYAN

uc_list = [
    "• Primary Actor: Placement Aspirant Candidate",
    "• Use Cases Executed:\n  1. Biometric Unlock App\n  2. Select Practice Category\n  3. Run JavaScript Code\n  4. Attempt Timed Exam\n  5. View Score & Rationale\n  6. Pin Revision Bookmarks"
]
for u in uc_list:
    p2 = tf.add_paragraph()
    p2.text = u
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(12)

s15_seq = [t15_id]
if uc_pic: s15_seq.append(uc_pic.shape_id)
s15_seq.append(card15.shape_id)
add_sequential_click_animations(slide15, s15_seq)

# SLIDE 16: SYSTEM IMPLEMENTATION
slide16 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide16)
t16_id = add_header(slide16, "System Implementation Highlights", "CODE & TECHNICAL LOGIC")

impl_items = [
    ("Native Android Shell (`MainActivity.kt`)", "Kotlin 1.9 & Jetpack Compose activity setting WindowManager `FLAG_SECURE` screenshot blocking."),
    ("Anti-Cheating Engine (`cheating_protection.js`)", "Monitors `window.blur` events and tab switches, triggering auto-submission upon focus loss."),
    ("JavaScript Code Runner (`coding_runner.js`)", "Client-side JS evaluator capturing stdout console logs and verifying PASS test case badges."),
    ("LocalStorage State Manager (`storage.js`)", "Centralized synchronous state dispatcher storing theme, bookmarks, and test history offline.")
]

s16_shapes = [t16_id]
for idx, (title, desc) in enumerate(impl_items):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide16, x, y, Inches(5.75), Inches(2.55), CARD_BORDER)
    s16_shapes.append(card.shape_id)
    
    tb = slide16.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.15), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)

add_sequential_click_animations(slide16, s16_shapes)

# SLIDE 17: UI SCREENSHOTS GALLERY
slide17 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide17)
t17_id = add_header(slide17, "Application UI Screenshots Gallery", "INTERFACE PREVIEWS")

snaps_to_show = [
    ("Dashboard UI", "snap_index.png"),
    ("Aptitude UI", "snap_aptitude.png"),
    ("Coding Sandbox UI", "snap_coding.png"),
    ("Timed Exam UI", "snap_mocktest.png")
]

s17_shapes = [t17_id]
for idx, (title, filename) in enumerate(snaps_to_show):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    sp = os.path.join(screenshots_dir, filename)
    if os.path.exists(sp):
        pic = slide17.shapes.add_picture(sp, x, y, Inches(5.75), Inches(2.6))
        s17_shapes.append(pic.shape_id)

add_sequential_click_animations(slide17, s17_shapes)

# SLIDE 18: SYSTEM TESTING & QA METRICS
slide18 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide18)
t18_id = add_header(slide18, "System Testing & Quality Assurance", "TESTING EXECUTION")

test_summary_items = [
    ("Testing Strategy Executed", "Unit, Integration, Security (FLAG_SECURE & Focus Loss Guard), Performance (<50ms load), and Responsiveness Testing."),
    ("20 Test Cases Matrix", "Executed 20 professional test cases covering Native Shell, Offline Storage, Coding Runner, and Mock Exam with 100% PASS rate."),
    ("Security & Integrity Check", "Confirmed 100% blocking of PrintScreen, Ctrl+C/V, F12 inspect keys, and right-click context menus during active exams.")
]

s18_shapes = [t18_id]
for idx, (title, desc) in enumerate(test_summary_items):
    y = Inches(1.3 + idx * 1.85)
    card = create_glass_card(slide18, Inches(0.7), y, Inches(11.933), Inches(1.68), CARD_BORDER)
    s18_shapes.append(card.shape_id)
    
    tb = slide18.shapes.add_textbox(Inches(1.0), y + Inches(0.18), Inches(11.3), Inches(1.35))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"✅  {title}"
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(6)

add_sequential_click_animations(slide18, s18_shapes)

# SLIDE 19: KEY ADVANTAGES & BENEFITS
slide19 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide19)
t19_id = add_header(slide19, "Key Advantages & System Benefits", "PLATFORM IMPACT")

advantages = [
    ("100% Offline-First Capability", "Operates under complete Airplane Mode with zero cloud server dependencies or data loss."),
    ("Sub-50ms DOM Rendering", "Instant local WebView rendering speed without ad banner delays or remote network latency."),
    ("Dual Security Protection", "Native Android Biometric fingerprint sensor lock paired with client-side window blur exam auto-submission."),
    ("$0 Operational Expenditure", "Serverless client-side execution and LocalStorage state management eliminating cloud infrastructure fees.")
]

s19_shapes = [t19_id]
for idx, (title, desc) in enumerate(advantages):
    col = idx % 2
    row = idx // 2
    x = Inches(0.7 + col * 5.95)
    y = Inches(1.3 + row * 2.75)
    card = create_glass_card(slide19, x, y, Inches(5.75), Inches(2.55), ACCENT_CYAN)
    s19_shapes.append(card.shape_id)
    
    tb = slide19.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.15), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"⭐ {title}"
    p.font.size = Pt(16.5)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)

add_sequential_click_animations(slide19, s19_shapes)

# SLIDE 20: FUTURE ENHANCEMENTS ROADMAP
slide20 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide20)
t20_id = add_header(slide20, "Future Enhancements Roadmap", "TECHNICAL INNOVATION")

roadmap_steps = [
    ("Phase 1", "Cloud Sync Middleware", "Optional Google Drive REST API integration for multi-device cloud backups."),
    ("Phase 2", "AI Diagnostic Analytics", "Automated weakness identification recommending targeted practice topics."),
    ("Phase 3", "WebSockets P2P Rooms", "Real-time multiplayer testing rooms for 1v1 coding battles."),
    ("Phase 4", "Instructor Portal", "Educator dashboard for building assessment profiles and exporting analytics.")
]

s20_shapes = [t20_id]
for idx, (phase, title, desc) in enumerate(roadmap_steps):
    x = Inches(0.7 + idx * 2.98)
    card = create_glass_card(slide20, x, Inches(1.3), Inches(2.8), Inches(5.4), ACCENT_INDIGO)
    s20_shapes.append(card.shape_id)
    
    tb = slide20.shapes.add_textbox(x + Inches(0.2), Inches(1.5), Inches(2.4), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = phase.upper()
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = HEADER_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(14)

add_sequential_click_animations(slide20, s20_shapes)

# SLIDE 21: CONCLUSION
slide21 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide21)
t21_id = add_header(slide21, "Conclusion & Key Achievements", "PROJECT SYNTHESIS")

card21 = create_glass_card(slide21, Inches(0.7), Inches(1.3), Inches(11.933), Inches(5.4), ACCENT_CYAN)
tb21 = slide21.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(4.8))
tf21 = tb21.text_frame
tf21.word_wrap = True

p = tf21.paragraphs[0]
p.text = "VetriPath Learn - Summary of Achievements:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = HEADER_CYAN

achievements = [
    "• Successfully designed and deployed an offline-first dual-platform application (Web PWA & Native Android APK).",
    "• Delivered zero-latency access to 2,000+ quantitative aptitude, reasoning, verbal, and coding practice problems.",
    "• Integrated client-side anti-cheating window blur focus loss monitoring and native Android FLAG_SECURE screenshot prevention.",
    "• Confirmed sub-50ms DOM rendering speed and 100% test case pass rates across unit, integration, and security matrices.",
    "• Operated at $0 operational expenditure using LocalStorage JSON state management, empowering students in low-bandwidth regions."
]

for a in achievements:
    p2 = tf21.add_paragraph()
    p2.text = a
    p2.font.size = Pt(14.5)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(14)

add_sequential_click_animations(slide21, [t21_id, card21.shape_id])

# SLIDE 22: THANK YOU WITH SEAMLESSLY MERGED CIRCULAR LOGO QR CODE
slide22 = prs.slides.add_slide(blank_layout)
add_slide_background_and_transition(slide22)

card22_left = create_glass_card(slide22, Inches(0.7), Inches(1.2), Inches(7.6), Inches(5.5), ACCENT_CYAN)
tb22 = slide22.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(7.2), Inches(4.7))
tf22 = tb22.text_frame
tf22.word_wrap = True

p = tf22.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "THANK YOU!"
p.font.size = Pt(46)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.font.name = "Segoe UI"

p2 = tf22.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.text = "VetriPath Learn — Offline-First Placement Platform"
p2.font.size = Pt(21)
p2.font.bold = True
p2.font.color.rgb = TEXT_WHITE
p2.space_before = Pt(14)

p3 = tf22.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.text = "\nQuestions & Technical Discussion"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = HEADER_CYAN
p3.space_before = Pt(14)

p4 = tf22.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
p4.text = "Presented by Darun Raj D & Nismetha P\nGuided by Dr. A. Meenakshi, M.C.A., M.Phil., M.S., Ph.D. (Assistant Professor)\nDept. of Computer Science & Applications"
p4.font.size = Pt(13)
p4.font.color.rgb = TEXT_MUTED
p4.space_before = Pt(12)

card22_right = create_glass_card(slide22, Inches(8.6), Inches(1.2), Inches(4.033), Inches(5.5), ACCENT_CYAN)

if os.path.exists(qr_path):
    slide22.shapes.add_picture(qr_path, Inches(9.116), Inches(1.5), width=Inches(3.0), height=Inches(3.0))

tb_qr = slide22.shapes.add_textbox(Inches(8.8), Inches(4.6), Inches(3.633), Inches(1.8))
tf_qr = tb_qr.text_frame
tf_qr.word_wrap = True

pq1 = tf_qr.paragraphs[0]
pq1.alignment = PP_ALIGN.CENTER
pq1.text = "SCAN TO VISIT LIVE PORTAL"
pq1.font.size = Pt(14.5)
pq1.font.bold = True
pq1.font.color.rgb = HEADER_CYAN

pq2 = tf_qr.add_paragraph()
pq2.alignment = PP_ALIGN.CENTER
pq2.text = "https://www.vetripathlearn.site/"
pq2.font.size = Pt(12.5)
pq2.font.bold = True
pq2.font.color.rgb = ACCENT_CYAN
pq2.space_before = Pt(4)

pq3 = tf_qr.add_paragraph()
pq3.alignment = PP_ALIGN.CENTER
pq3.text = "Scan with mobile camera to launch platform"
pq3.font.size = Pt(11)
pq3.font.color.rgb = TEXT_MUTED
pq3.space_before = Pt(4)

add_sequential_click_animations(slide22, [card22_left.shape_id, card22_right.shape_id])

# Save Presentation
prs.save(output_ppt)
prs.save(output_ppt_alt)
print(f"Presentation generated with professional animation timing: {output_ppt}")
