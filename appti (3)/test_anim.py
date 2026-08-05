import pptx
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Slide Transition (Fade Transition)
trans_xml = parse_xml('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>')
slide.element.append(trans_xml)

# Create 3 test shapes
shape1 = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 1000000, 1000000, 2000000, 1000000)
shape2 = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 1000000, 2500000, 2000000, 1000000)
shape3 = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 1000000, 4000000, 2000000, 1000000)

shape1.name = "Shape1"
shape2.name = "Shape2"
shape3.name = "Shape3"

# Function to add sequential click entrance animations (Appear / Fade In on click)
def add_sequential_entrance_animation(slide, shape_ids):
    timing_xml_str = f'''
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="root">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
    '''
    
    ctn_id = 3
    for idx, shape_id in enumerate(shape_ids):
        timing_xml_str += f'''
                    <p:par>
                      <p:cTn id="{ctn_id}" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="{ctn_id+1}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:set>
                                  <p:cb>
                                    <p:subSpTree>
                                      <p:cTn id="{ctn_id+2}" dur="1" fill="hold">
                                        <p:stCondLst>
                                          <p:cond delay="0"/>
                                        </p:stCondLst>
                                        <p:childTnLst>
                                          <p:set>
                                            <p:cb>
                                              <p:spTarget spid="{shape_id}"/>
                                            </p:cb>
                                            <p:to>
                                              <p:visibility val="visible"/>
                                            </p:to>
                                          </p:set>
                                        </p:childTnLst>
                                      </p:cTn>
                                    </p:subSpTree>
                                  </p:cb>
                                  <p:to>
                                    <p:style sub="fill"/>
                                  </p:to>
                                </p:set>
                                <p:anim filter="blinds(horizontal)" calcmode="lin" valueType="num">
                                  <p:cb>
                                    <p:spTarget spid="{shape_id}"/>
                                  </p:cb>
                                </p:anim>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
        '''
        ctn_id += 3

    timing_xml_str += '''
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    '''
    
    slide.element.append(parse_xml(timing_xml_str))

add_sequential_entrance_animation(slide, [shape1.shape_id, shape2.shape_id, shape3.shape_id])
prs.save("test_anim.pptx")
print("Saved presentation with animations!")
