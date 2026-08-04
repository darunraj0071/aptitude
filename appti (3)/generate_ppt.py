import os
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Paths & Logo Setup
logo_path = r"c:\Users\sanja\Downloads\appti (2) (1)\appti (3)\appti\images\logo.png"
screenshots_dir = r"c:\Users\sanja\Downloads\appti (2) (1)\appti (3)\screenshots"
img_dir = r"c:\Users\sanja\Downloads\appti (2) (1)\appti (3)\appti\images"
faded_logo_path = "faded_logo.png"

dfd_img = os.path.join(img_dir, "dfd_diagram.png")
er_img = os.path.join(img_dir, "er_diagram.png")
arch_img = os.path.join(img_dir, "architecture_diagram.png")

try:
    logo = Image.open(logo_path).convert("RGBA")
    r, g, b, alpha = logo.split()
    alpha_faded = ImageEnhance.Brightness(alpha).enhance(0.08)
    faded_logo = Image.merge("RGBA", (r, g, b, alpha_faded))
    faded_logo.save(faded_logo_path)
    print("Faded logo watermark created successfully.")
except Exception as e:
    print(f"Error creating faded logo: {e}")
    if os.path.exists(logo_path):
        import shutil
        shutil.copy(logo_path, faded_logo_path)

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_COLOR = RGBColor(15, 14, 23)        # #0F0E17
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_SILVER = RGBColor(148, 161, 178)
ACCENT_PURPLE = RGBColor(127, 90, 240)
ACCENT_GREEN = RGBColor(44, 182, 125)
ACCENT_BLUE = RGBColor(56, 189, 248)
PANEL_BG = RGBColor(31, 30, 39)

def apply_slide_bg(slide, use_watermark=True, add_header_logo=True):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if use_watermark and os.path.exists(faded_logo_path):
        width_in = 5.5
        height_in = 5.5
        left_in = (13.333 - width_in) / 2
        top_in = (7.5 - height_in) / 2
        slide.shapes.add_picture(faded_logo_path, Inches(left_in), Inches(top_in), width=Inches(width_in), height=Inches(height_in))

    if add_header_logo and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(12.2), Inches(0.3), width=Inches(0.6), height=Inches(0.6))

def add_slide_title(slide, text, category=None):
    if category:
        cat_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.2), Inches(11.0), Inches(0.3))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(10.5)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_GREEN
        p_c.font.name = "Segoe UI"

    top_pos = Inches(0.48) if category else Inches(0.35)
    title_box = slide.shapes.add_textbox(Inches(0.7), top_pos, Inches(11.0), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = TEXT_WHITE
    
    line_top = top_pos + Inches(0.52)
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), line_top, Inches(2.2), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ACCENT_PURPLE
    underline.line.fill.background()

def add_bullet_slide(prs, title, bullets, category=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, use_watermark=True, add_header_logo=True)
    add_slide_title(slide, title, category=category)
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.93), Inches(5.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.level = 0
        
        if isinstance(item, tuple):
            lead, body = item
            run = p.add_run()
            run.text = "• " + lead + ": "
            run.font.bold = True
            run.font.size = Pt(15)
            run.font.color.rgb = ACCENT_PURPLE
            run.font.name = "Segoe UI"
            
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(15)
            run2.font.color.rgb = TEXT_WHITE
            run2.font.name = "Segoe UI"
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(15)
            run.font.color.rgb = TEXT_WHITE
            run.font.name = "Segoe UI"
    return slide

def add_image_slide(prs, title, img_path, description, category=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, use_watermark=True, add_header_logo=True)
    add_slide_title(slide, title, category=category)

    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.93), Inches(0.5))
        p = desc_box.text_frame.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_SILVER
        p.font.name = "Segoe UI"

    if os.path.exists(img_path):
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.7), Inches(10.93), Inches(5.3))
        frame.fill.solid()
        frame.fill.fore_color.rgb = PANEL_BG
        frame.line.color.rgb = ACCENT_GREEN
        frame.line.width = Pt(1.5)

        slide.shapes.add_picture(img_path, Inches(1.35), Inches(1.85), width=Inches(10.63), height=Inches(5.0))
    return slide

def add_page_io_screenshot_slide(prs, page_title, page_url, snap_file, overview_text, input_points, output_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, use_watermark=True, add_header_logo=True)
    add_slide_title(slide, page_title, category=f"SYSTEM DESIGN - MODULE SPECIFICATION ({page_url})")
    
    top_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.15), Inches(11.93), Inches(0.75))
    top_card.fill.solid()
    top_card.fill.fore_color.rgb = PANEL_BG
    top_card.line.color.rgb = ACCENT_PURPLE
    top_card.line.width = Pt(1.5)
    
    tf_top = top_card.text_frame
    tf_top.word_wrap = True
    tf_top.margin_left = tf_top.margin_right = Inches(0.18)
    tf_top.margin_top = Inches(0.08)
    p_t1 = tf_top.paragraphs[0]
    p_t1.text = "Overview & Purpose: " + overview_text
    p_t1.font.size = Pt(11)
    p_t1.font.color.rgb = TEXT_WHITE
    p_t1.font.name = "Segoe UI"

    io_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.0), Inches(6.0), Inches(5.1))
    io_card.fill.solid()
    io_card.fill.fore_color.rgb = PANEL_BG
    io_card.line.color.rgb = ACCENT_BLUE
    io_card.line.width = Pt(1.5)
    
    tf_io = io_card.text_frame
    tf_io.word_wrap = True
    tf_io.margin_left = tf_io.margin_right = Inches(0.2)
    tf_io.margin_top = Inches(0.15)

    p_in_h = tf_io.paragraphs[0]
    p_in_h.text = "📥 USER INPUT SPECIFICATIONS"
    p_in_h.font.size = Pt(12)
    p_in_h.font.bold = True
    p_in_h.font.color.rgb = ACCENT_BLUE
    p_in_h.font.name = "Segoe UI"
    p_in_h.space_after = Pt(4)
    
    for inp in input_points:
        p = tf_io.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(inp, tuple):
            label, desc = inp
            r1 = p.add_run()
            r1.text = "• " + label + ": "
            r1.font.bold = True
            r1.font.size = Pt(10)
            r1.font.color.rgb = TEXT_WHITE
            r1.font.name = "Segoe UI"
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = TEXT_SILVER
            r2.font.name = "Segoe UI"

    p_out_h = tf_io.add_paragraph()
    p_out_h.space_before = Pt(6)
    p_out_h.text = "📤 SYSTEM OUTPUT SPECIFICATIONS"
    p_out_h.font.size = Pt(12)
    p_out_h.font.bold = True
    p_out_h.font.color.rgb = ACCENT_GREEN
    p_out_h.font.name = "Segoe UI"
    p_out_h.space_after = Pt(4)
    
    for outp in output_points:
        p = tf_io.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(outp, tuple):
            label, desc = outp
            r1 = p.add_run()
            r1.text = "• " + label + ": "
            r1.font.bold = True
            r1.font.size = Pt(10)
            r1.font.color.rgb = TEXT_WHITE
            r1.font.name = "Segoe UI"
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = TEXT_SILVER
            r2.font.name = "Segoe UI"

    snap_path = os.path.join(screenshots_dir, snap_file)
    if os.path.exists(snap_path):
        img_frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.73), Inches(5.1))
        img_frame.fill.solid()
        img_frame.fill.fore_color.rgb = PANEL_BG
        img_frame.line.color.rgb = ACCENT_GREEN
        img_frame.line.width = Pt(1.5)
        
        tf_fr = img_frame.text_frame
        tf_fr.margin_top = Inches(0.08)
        p_fr = tf_fr.paragraphs[0]
        p_fr.text = f"📱 REAL UI SCREENSHOT: {page_url}"
        p_fr.font.size = Pt(11)
        p_fr.font.bold = True
        p_fr.font.color.rgb = ACCENT_GREEN
        p_fr.font.name = "Segoe UI"
        p_fr.alignment = PP_ALIGN.CENTER
        
        slide.shapes.add_picture(snap_path, Inches(7.02), Inches(2.45), width=Inches(5.49), height=Inches(4.5))

    return slide

# ==================== SLIDE PRESENTATION BUILDER ====================

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide1, use_watermark=False, add_header_logo=False)
if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(5.41), Inches(0.8), width=Inches(2.5), height=Inches(2.5))

tx1 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.33), Inches(3.5))
tf1 = tx1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "VETRIPATH"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.font.name = "Segoe UI"
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(8)

p2 = tf1.add_paragraph()
p2.text = "Complete Offline Placement Preparation Portal & Native Mobile Application"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = ACCENT_PURPLE
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER
p2.space_after = Pt(16)

p3 = tf1.add_paragraph()
p3.text = "Academic Technical Presentation • Department of Computer Science & Engineering"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_SILVER
p3.font.name = "Segoe UI"
p3.alignment = PP_ALIGN.CENTER

# Slide 2: Abstract
add_bullet_slide(prs, "PROJECT ABSTRACT", [
    ("Concept & Mission", "VetriPath is a self-contained hybrid education and assessment portal designed to prepare students for campus placements, technical interviews, and competitive exams."),
    ("Offline-First Engineering", "Eliminates network latency by embedding the web platform (HTML5/CSS3/ES6 JS) directly within an Android native WebView container, running 100% offline."),
    ("Multi-Modal Learning Modules", "Provides Quantitative Aptitude, Logical Reasoning, Verbal Ability, an Interactive JavaScript Coding Sandbox, and placement Roadmaps."),
    ("Exam Integrity Suite", "Features client-side anti-cheating mechanisms including shortcut suppression, copy-paste blocking, and instant exam auto-submission upon window tab blur."),
    ("Native Biometric Shield", "Employs Jetpack Compose and Android Biometrics SDK to secure local performance logs and progress history.")
], category="ABSTRACT")

# Slide 3: SYSTEM ANALYSIS - Feasibility Study
add_bullet_slide(prs, "SYSTEM ANALYSIS: Feasibility Study", [
    ("Technical Feasibility", "Built on standardized web engines (HTML5/CSS3/ES6 JS) & native WebView interop. Zero network requests, 100% offline operational stability."),
    ("Economic Feasibility", "Serverless client-side architecture with zero database or server hosting fees. Developed using 100% open-source tools ($0 overhead)."),
    ("Operational Feasibility", "Intuitive glassmorphic UI design requiring zero user training. Automated scoring, feedback, and history management ensure high adoption.")
], category="SYSTEM ANALYSIS")

# Slide 4: SYSTEM ANALYSIS - Existing System & Drawbacks
add_bullet_slide(prs, "SYSTEM ANALYSIS: Existing System & Drawbacks", [
    ("Network Dependencies", "Online portals (GeeksforGeeks, IndiaBIX, HackerRank) fail completely in low-bandwidth or rural regions."),
    ("High Ads & Latency", "Heavy ad banner networks slow down page rendering, consume high RAM, and distract study focus."),
    ("Insecure Practice Sandbox", "Traditional portals permit copy-pasting code, viewing Developer Console inspector (F12), and changing tabs during mock tests.")
], category="SYSTEM ANALYSIS")

# Slide 5: SYSTEM ANALYSIS - Proposed System & Benefits
add_bullet_slide(prs, "SYSTEM ANALYSIS: Proposed System & Benefits", [
    ("100% Offline Asset Packaging", "Direct asset integration inside Android APK assets directory. Zero loading latencies."),
    ("Jetpack Compose Biometric Shield", "Protects user data files using Android biometrics (Fingerprint/PIN prompts) prior to app load."),
    ("Cheating Protection Suite", "Monitors keyboard events & auto-submits mock exams immediately when tab switch or window blur occurs.")
], category="SYSTEM ANALYSIS")

# Slide 6: SYSTEM ANALYSIS - Scope of the Project
add_bullet_slide(prs, "SYSTEM ANALYSIS: Scope of the Project", [
    ("Target Audience", "Undergraduate & postgraduate students preparing for campus drives (TCS NQT, Infosys, Wipro, Accenture, CTS, product startups)."),
    ("Module Coverage", "Quantitative Aptitude (14 topics), Logical Reasoning (9 topics), Verbal Ability (7 topics), JS Coding Sandbox, Timed Mock Exams, Roadmaps."),
    ("Platform Support", "Cross-platform accessibility across mobile viewports, tablet screens, and desktop web environments.")
], category="SYSTEM ANALYSIS")

# Slide 7: SYSTEM SPECIFICATION - Hardware & Software Configuration
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide7, use_watermark=True, add_header_logo=True)
add_slide_title(slide7, "SYSTEM SPECIFICATION: Hardware & Software", category="SYSTEM SPECIFICATION")

sw_bullets = [
    ("Operating System", "Android 8.0 (API 26) or higher for APK; Windows 10/11 for Web"),
    ("Core Languages", "Kotlin 1.9, Java (Android SDK), JavaScript (ES6+), HTML5, CSS3"),
    ("Frameworks & SDKs", "Jetpack Compose, Android Jetpack Biometrics SDK, WebView Engine"),
    ("Tools & Compilers", "Android Studio Hedgehog, VS Code, Gradle Build Tool, Python 3.10+")
]
hw_bullets = [
    ("Development Unit", "Intel Core i5 / AMD Ryzen 5, 8 GB RAM (16 GB recommended), 10 GB storage"),
    ("Testing Device", "Android Smartphone with Biometric Sensor (Fingerprint/PIN), Touch screen interface"),
    ("Client Display", "Desktop resolution 1366x768 pixels; mobile viewports min 360x640 pixels")
]

tx1 = slide7.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.7), Inches(5.5))
tf1 = tx1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "Software Configuration"
p1.font.size = Pt(17)
p1.font.bold = True
p1.font.color.rgb = ACCENT_GREEN
p1.font.name = "Segoe UI"
p1.space_after = Pt(10)

for item in sw_bullets:
    p = tf1.add_paragraph()
    p.space_after = Pt(8)
    lead, body = item
    r1 = p.add_run()
    r1.text = "• " + lead + ": "
    r1.font.bold = True
    r1.font.size = Pt(12.5)
    r1.font.color.rgb = ACCENT_PURPLE
    r1.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Segoe UI"

tx2 = slide7.shapes.add_textbox(Inches(6.9), Inches(1.35), Inches(5.7), Inches(5.5))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Hardware Configuration"
p2.font.size = Pt(17)
p2.font.bold = True
p2.font.color.rgb = ACCENT_GREEN
p2.font.name = "Segoe UI"
p2.space_after = Pt(10)

for item in hw_bullets:
    p = tf2.add_paragraph()
    p.space_after = Pt(8)
    lead, body = item
    r1 = p.add_run()
    r1.text = "• " + lead + ": "
    r1.font.bold = True
    r1.font.size = Pt(12.5)
    r1.font.color.rgb = ACCENT_PURPLE
    r1.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Segoe UI"

# ==================== ALL 12 WEB UI SCREENSHOT SLIDES ====================

web_modules_slides = [
    ("Dashboard UI Console", "index.html", "snap_index.png", "Main control center rendering overall progress, total available questions (2000+), solved metrics, and action cards.", [
        ("Navigation Clicks", "Selects module links from header or drawer."),
        ("Quick Action CTAs", "Clicks 'Start Practice' or 'Take Mock Exam' cards."),
        ("Theme Mode Switch", "Toggles light/dark interface style icon button.")
    ], [
        ("Metric Cards", "Displays available questions count (2000+) and solved metrics."),
        ("Progress Bar", "Visual completion percentage meter from LocalStorage."),
        ("Feature Tiles", "Interactive topic tiles with direct routing.")
    ]),
    ("Quantitative Aptitude Module", "aptitude.html", "snap_aptitude.png", "Provides mathematical problem-solving practice across 14 core topics with formula Guides.", [
        ("Topic Selection", "Selects mathematical topic from category drawer."),
        ("Option Radio Click", "Clicks option choice (A, B, C, or D) for a question."),
        ("Solution Toggle", "Clicks 'Show Solution' button to view formula breakdown.")
    ], [
        ("Formula Reference", "Displays mathematical formulas at top of topic."),
        ("Instant Verdict", "Highlights chosen option green (Correct) or red (Incorrect)."),
        ("Step-by-Step Rationale", "Shows detailed step-by-step mathematical breakdown.")
    ]),
    ("Logical Reasoning Module", "reasoning.html", "snap_reasoning.png", "Focuses on analytical ability, pattern identification, and deduction across 9 topics.", [
        ("Reasoning Category", "Chooses topic (Blood Relations, Syllogisms, Coding-Decoding)."),
        ("Answer Selection", "Clicks answer choice matching logical deduction."),
        ("Rule Sheet Toggle", "Clicks 'View Reasoning Rules' button.")
    ], [
        ("Rule Box", "Renders analytical guidelines and family tree rules."),
        ("Immediate Validation", "Highlights correct answer choice and marks wrong attempt."),
        ("Logic Breakdown", "Displays step-by-step logical deduction explanation.")
    ]),
    ("Verbal Ability Module", "verbal.html", "snap_verbal.png", "Enhances English vocabulary, grammar accuracy, and comprehension skills across 7 topics.", [
        ("Sub-category Choice", "Selects verbal module (Grammar Rules, Error Spotting)."),
        ("Choice Selection", "Clicks selected multiple-choice answer option."),
        ("Passage Navigation", "Scrolls and reads comprehension text passages.")
    ], [
        ("Grammar Tip Card", "Presents essential English usage rules before questions."),
        ("Answer Feedback", "Shows immediate correctness indicator with explanation."),
        ("Contextual Rationale", "Explains why specific option is grammatically sound.")
    ]),
    ("Interactive Coding Workstation", "coding.html", "snap_coding.png", "Client-side JavaScript execution sandbox allowing candidates to write and test code.", [
        ("Problem Selector", "Picks problem (Two Sum, Palindrome Check, Array Reverse)."),
        ("Code Editor Input", "Writes or edits JavaScript code inside interactive text area."),
        ("Run Code Trigger", "Clicks 'Execute Code' trigger to compile script.")
    ], [
        ("Console Output (stdout)", "Captures and displays console.log output and return values."),
        ("Test Case Badges", "Displays PASS / FAIL status for default test inputs."),
        ("Runtime Error Log", "Displays syntax error messages or stack trace logs.")
    ]),
    ("Practice Hub Topic Directory", "practice.html", "snap_practice.png", "Central directory listing all topics with multi-category tabs and search filters.", [
        ("Category Filter Tabs", "Clicks 'All', 'Quantitative', 'Reasoning', 'Verbal', or 'Coding' tabs."),
        ("Search Keyword Input", "Types topic title or keyword into search bar."),
        ("Topic Card Click", "Clicks topic card to launch practice session.")
    ], [
        ("Filtered Grid", "Renders topic cards matching selected tab and search query."),
        ("Completion Badges", "Shows completed question count vs total questions for each topic."),
        ("Action CTAs", "Provides 'Continue' or 'Start Fresh' buttons.")
    ]),
    ("Mock Test Exam Simulator", "mocktest.html", "snap_mocktest.png", "Timed corporate exam simulator with randomized question pools and tab blur auto-submit.", [
        ("Tier Setup", "Selects tier: Easy (15m), Medium (30m), Hard (45m), Full-Length."),
        ("Option Selection", "Selects radio buttons for questions under timed environment."),
        ("Tab Blur / Switch", "Candidate switches browser tab or minimizes window during test.")
    ], [
        ("Live Timer & Grid", "Renders ticking countdown clock and question state palette."),
        ("Security Modal & Auto-Submit", "Triggers red security modal & auto-submits exam upon focus loss!"),
        ("Result Analytics", "Displays final Score, Accuracy %, Time Taken, and detailed breakdown.")
    ]),
    ("Bookmarks Review Hub", "bookmarks.html", "snap_bookmarks.png", "Dedicated review repository maintaining pinned questions across modules.", [
        ("Category Filter", "Filters bookmarked questions by Subject."),
        ("Remove Bookmark", "Clicks trash icon on question card to unpin."),
        ("Clear All Button", "Clicks 'Clear All Bookmarks' with confirmation modal.")
    ], [
        ("Saved Question Cards", "Displays full question prompt, choices, and stored correct answer."),
        ("Empty State Graphic", "Shows friendly illustration when no bookmarks are saved."),
        ("Updated Local State", "Removes unpinned questions from LocalStorage array instantly.")
    ]),
    ("Placement Preparation Roadmap", "roadmap.html", "snap_roadmap.png", "Structured 6-month preparation schedule checklist for campus placements.", [
        ("Track Filter", "User chooses target profile (Product Companies, Service MNCs, Core)."),
        ("Milestone Checkbox", "Checks off completed preparation tasks and topics.")
    ], [
        ("Timeline Grid", "Renders 6-month roadmap stages from foundational math to mock interviews."),
        ("Readiness Bar", "Calculates overall career readiness progress percentage."),
        ("Checklist State Sync", "Persists checked milestone status in LocalStorage.")
    ]),
    ("Global Search Index UI", "search.html", "snap_search.png", "Enables real-time keyword search indexing across all 30+ topics.", [
        ("Search Input Text", "Types search term (e.g. 'Percentage', 'Trees', 'Syllogism', 'Arrays')."),
        ("Filter Button", "Clicks clear or topic category quick tag.")
    ], [
        ("Real-time Results", "Instantly updates search results showing topic title and direct link."),
        ("Zero Results Slate", "Suggests alternative search terms if no matching topic is found.")
    ]),
    ("About VetriPath & Tech Stack", "about.html", "snap_about.png", "Presents platform architecture, offline-first engineering principles, and tech stack.", [
        ("Theme Mode Switch", "Clicks Dark Mode / Light Mode toggle icon."),
        ("Feedback Link", "Clicks developer contact or app review button.")
    ], [
        ("Architecture Summary", "Highlights offline WebView benefits and biometric security integration."),
        ("Tech Stack Badges", "Displays HTML5, CSS3, JS ES6, Kotlin, Compose, and WebSockets badges."),
        ("Global Theme Sync", "Applies theme preference across all web pages.")
    ]),
    ("Contact & Feedback Form", "contact.html", "snap_contact.png", "Provides an interactive user feedback form for reporting issues and feature requests.", [
        ("User Input Form", "User inputs Name, Email ID, Feedback Category, and Message text."),
        ("Submit Button Click", "Clicks 'Submit Feedback' form button.")
    ], [
        ("Form Validation", "Checks required fields and email formatting."),
        ("Success Toast", "Displays animated green toast: 'Feedback saved successfully!'"),
        ("Local Feedback Log", "Saves feedback entry locally in browser storage.")
    ])
]

for title, url, snap, overview, inputs, outputs in web_modules_slides:
    add_page_io_screenshot_slide(
        prs,
        page_title=title,
        page_url=url,
        snap_file=snap,
        overview_text=overview,
        input_points=inputs,
        output_points=outputs
    )

# ==================== ALL 5 ACTIVE QUIZ ATTENDING SAMPLE SLIDES ====================

active_quiz_slides = [
    ("Quantitative Aptitude Active Practice Workstation", "quiz_sample_aptitude.png", "Candidate attending Percentages question, radio choice B selected, green correct answer feedback & formula breakdown visible.", "ACTIVE QUIZ ATTENDING"),
    ("Logical Reasoning Active Practice Workstation", "quiz_sample_reasoning.png", "Candidate attending Blood Relations question, analytical family tree symbols, radio option selected & logic tree visible.", "ACTIVE QUIZ ATTENDING"),
    ("Verbal Ability Active Practice Workstation", "quiz_sample_verbal.png", "Candidate attending Grammar Rules question, sentence option choice selected & subject-verb rationale expanded.", "ACTIVE QUIZ ATTENDING"),
    ("Interactive Coding Sandbox Active Execution Sample", "quiz_sample_coding.png", "Candidate editing Two Sum JavaScript solution, 'Execute Code' triggered, PASS/FAIL test badges & stdout console output visible.", "ACTIVE QUIZ ATTENDING"),
    ("Timed Exam Simulator Active Test Workstation", "quiz_sample_mocktest.png", "Candidate taking 30-min placement mock exam, live countdown clock running, question palette grid active & tab blur monitor engaged.", "ACTIVE QUIZ ATTENDING")
]

for title, snap, desc, cat in active_quiz_slides:
    s_path = os.path.join(screenshots_dir, snap)
    add_image_slide(prs, title, s_path, desc, category=cat)

# ==================== SYSTEM DESIGN DIAGRAM SLIDES ====================
add_image_slide(prs, "SYSTEM DESIGN: Level 1 Dataflow Diagram (DFD)", dfd_img, "Level 1 Dataflow Diagram mapping candidate inputs through processes to LocalStorage.", category="SYSTEM DESIGN")
add_image_slide(prs, "SYSTEM DESIGN: Hybrid System Architecture Diagram", arch_img, "4-Layer Hybrid Native Architecture: UI Layer -> Compose Shell -> WebView Host -> LocalStorage DB.", category="SYSTEM DESIGN")
add_image_slide(prs, "SYSTEM DESIGN: Database Schema (ER Diagram)", er_img, "LocalStorage Entity-Relationship Diagram detailing JSON schema keys and parameters.", category="SYSTEM DESIGN")

# ==================== SYSTEM IMPLEMENTATION SLIDES ====================
add_bullet_slide(prs, "SYSTEM IMPLEMENTATION: Coding Highlights", [
    ("Native Kotlin Security Shell", "Applies FLAG_SECURE to prevent screenshots/recording and initializes BiometricPrompt scanner."),
    ("JavaScript State Storage (storage.js)", "Consolidates getState(), updateState(), addBookmark(), and logSolvedQuestion() state hooks."),
    ("Exam Anti-Cheating Suite", "Monitors window blur and visibility state changes, auto-submitting exams upon focus loss.")
], category="SYSTEM IMPLEMENTATION")

# ==================== SYSTEM TESTING SLIDES ====================
add_bullet_slide(prs, "SYSTEM TESTING: Test Results Summary", [
    ("Biometric Authentication Lock", "TC-01: Verified BiometricPrompt scanner overlays and PIN fallbacks [STATUS: PASS]."),
    ("Offline Airplane Mode Load", "TC-02: Confirmed asset loading <50ms with zero network requests [STATUS: PASS]."),
    ("Exam Blur Auto-Submission", "TC-03: Simulated window blur during active exam; 100% detection rate [STATUS: PASS]."),
    ("Shortcut & Context Menu Disable", "TC-04: Confirmed suppression of Ctrl+C, Ctrl+V, Ctrl+U, F12 inspect keys [STATUS: PASS].")
], category="SYSTEM TESTING")

# ==================== FUTURE ENHANCEMENTS & CONCLUSION SLIDES ====================
add_bullet_slide(prs, "FUTURE ENHANCEMENTS & CONCLUSION", [
    ("Cloud Sync Middleware", "Integration with Google Drive REST API for optional cloud backup and restore."),
    ("AI Personal Analytics", "Recommends targeted prep paths based on solved question difficulty scores."),
    ("Peer-to-Peer Competitive Rooms", "Online testing rooms and real-time scoreboards using WebSockets."),
    ("Conclusion", "Successfully developed a high-performance offline hybrid platform bridging the digital divide for students.")
], category="FUTURE ENHANCEMENTS")

# Thank You Slide
slide20 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide20, use_watermark=False, add_header_logo=False)
if os.path.exists(logo_path):
    slide20.shapes.add_picture(logo_path, Inches(5.41), Inches(1.2), width=Inches(2.5), height=Inches(2.5))

tx20 = slide20.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(3.0))
tf20 = tx20.text_frame
tf20.word_wrap = True

p_ty = tf20.paragraphs[0]
p_ty.text = "THANK YOU"
p_ty.font.size = Pt(54)
p_ty.font.bold = True
p_ty.font.color.rgb = TEXT_WHITE
p_ty.font.name = "Segoe UI"
p_ty.alignment = PP_ALIGN.CENTER
p_ty.space_after = Pt(12)

p_tag = tf20.add_paragraph()
p_tag.text = '"Walk the Path. Become the Victory. Inspire the World."'
p_tag.font.size = Pt(20)
p_tag.font.italic = True
p_tag.font.color.rgb = ACCENT_GREEN
p_tag.font.name = "Segoe UI"
p_tag.alignment = PP_ALIGN.CENTER

output_name = "VetriPath_Project_Presentation.pptx"
prs.save(output_name)
print(f"Presentation with ALL 20 IMAGES saved successfully as {output_name}.")

if os.path.exists(faded_logo_path):
    os.remove(faded_logo_path)
    print("Cleaned up temporary files.")
